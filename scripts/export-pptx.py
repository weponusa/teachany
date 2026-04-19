#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TeachAny · HTML → PPTX 导出工具（v5.34）

用法:
    python3 scripts/export-pptx.py <课件目录>
    python3 scripts/export-pptx.py examples/math-m-vertex-form

行为:
    - 读取 <课件目录>/index.html + <课件目录>/manifest.json
    - 按 <section> 切分幻灯片 (Hero / 学习目标 / 前测 / 每个知识模块 / 综合任务 / 后测 / 小结)
    - 自动抓取 <img src="./assets/*.png"> 作为幻灯片主图
    - 文本保留标题层级 (h1→大标题 / h2→二级标题 / p→正文)
    - 练习题保留题干 + 正确答案 (供教师使用)
    - 互动组件 (canvas / 知识图谱 / 音频播放器 / AI 学伴) 在 PPTX 版降级为
      "线上互动入口" 占位页 + URL 链接
    - 产出: <课件目录>/<课件名>.pptx

依赖:
    pip3 install python-pptx beautifulsoup4

遵循硬规则 #46 PPTX 导出基线 (TeachAny SKILL_CN.md v5.34)。
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

# ───────────────────────────── 依赖自动安装 ─────────────────────────────

def _ensure_deps():
    missing = []
    try:
        import pptx  # noqa: F401
    except ImportError:
        missing.append('python-pptx')
    try:
        from bs4 import BeautifulSoup  # noqa: F401
    except ImportError:
        missing.append('beautifulsoup4')

    if missing:
        print(f"🔧 检测到缺失依赖: {', '.join(missing)}，尝试自动安装...")
        import subprocess
        ret = subprocess.call(
            [sys.executable, '-m', 'pip', 'install', '--quiet', *missing]
        )
        if ret != 0:
            print(f"❌ 自动安装失败。请手动执行: pip3 install {' '.join(missing)}")
            sys.exit(2)
        print("✅ 依赖安装完成")


_ensure_deps()

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from bs4 import BeautifulSoup  # noqa: E402


# ───────────────────────────── 常量 ─────────────────────────────

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)

PALETTE = {
    'primary': RGBColor(0x63, 0x66, 0xf1),
    'primary_dark': RGBColor(0x4f, 0x46, 0xe5),
    'accent': RGBColor(0xec, 0x48, 0x99),
    'text_dark': RGBColor(0x1e, 0x29, 0x3b),
    'text_soft': RGBColor(0x64, 0x74, 0x8b),
    'bg_soft': RGBColor(0xf8, 0xfa, 0xfc),
    'placeholder': RGBColor(0xee, 0xf2, 0xff),
}

# 互动组件关键词 → 占位页
INTERACTIVE_MARKERS = [
    ('canvas', '🎨 Canvas 互动'),
    ('#knowledge-graph', '🗺️ 知识图谱'),
    ('audioPlaylist', '🔊 语音讲解（滚动自动播放）'),
    ('ai-tutor', '🤖 AI 学伴悬浮球'),
    ('remotion', '🎬 Remotion 教学动画'),
]


# ───────────────────────────── 工具函数 ─────────────────────────────

def clean_text(s: str) -> str:
    """清理多余空白 + 删零宽字符"""
    if not s:
        return ''
    s = re.sub(r'\s+', ' ', s)
    s = s.replace('\u200b', '').replace('\ufeff', '').strip()
    return s


def read_manifest(course_dir: Path) -> dict:
    p = course_dir / 'manifest.json'
    if not p.exists():
        return {}
    try:
        return json.loads(p.read_text(encoding='utf-8'))
    except Exception as e:
        print(f"⚠️  manifest 解析失败: {e}")
        return {}


def find_first_image(section, course_dir: Path):
    """从 section 中找第一张存在的本地图片，返回绝对路径 Path 或 None"""
    for img in section.find_all('img'):
        src = img.get('src', '').strip()
        if not src or src.startswith('data:'):
            continue
        if src.startswith('http://') or src.startswith('https://'):
            continue
        candidate = (course_dir / src).resolve()
        # 在 course_dir 范围内才接受
        try:
            candidate.relative_to(course_dir.resolve())
        except ValueError:
            continue
        if candidate.exists() and candidate.suffix.lower() in ('.png', '.jpg', '.jpeg', '.gif', '.webp'):
            return candidate
    return None


def detect_interactive_components(section) -> list:
    """检查 section 内是否有互动组件，返回要在 PPTX 中降级提示的组件清单"""
    html_str = str(section).lower()
    found = []
    for marker, label in INTERACTIVE_MARKERS:
        if marker.lower() in html_str:
            found.append(label)
    return found


def extract_section_texts(section) -> dict:
    """从 section 中抽出结构化文本"""
    # 标题优先级 h1 > h2 > h3
    title = ''
    for tag in ['h1', 'h2', 'h3']:
        h = section.find(tag)
        if h:
            title = clean_text(h.get_text())
            break

    # 二级小标题 & 段落
    bullets = []
    for el in section.find_all(['h3', 'h4', 'p', 'li']):
        txt = clean_text(el.get_text())
        if not txt or txt == title:
            continue
        # 去重
        if bullets and bullets[-1] == txt:
            continue
        # 过滤极短无意义文本
        if len(txt) < 2:
            continue
        bullets.append(txt)

    # 练习题: 找 .quiz, .question, [data-quiz-id] 或 button onclick="handleQuiz(..."
    quizzes = []
    for quiz_el in section.select('.quiz, .question, [data-quiz-id]'):
        q_text = clean_text(quiz_el.get_text())
        if q_text:
            quizzes.append(q_text[:500])  # 裁长
    # 从 onclick 提取答案提示
    for btn in section.find_all(attrs={'onclick': True}):
        onclick = btn.get('onclick', '')
        m = re.search(r"handleQuiz\s*\(\s*this\s*,\s*['\"]([^'\"]+)['\"]\s*,\s*['\"]([A-Z])['\"]", onclick)
        if m:
            qid, correct = m.group(1), m.group(2)
            if not any(f'[{qid}]' in q for q in quizzes):
                quizzes.append(f'[{qid}] 正确答案: {correct}')

    return {
        'title': title,
        'bullets': bullets[:12],   # 避免单页过长
        'quizzes': quizzes[:4],
    }


# ───────────────────────────── PPTX 构建 ─────────────────────────────

def setup_pres():
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    return pres


def add_title_slide(pres, course_title: str, subtitle: str):
    """首页（Hero）"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank

    # 渐变背景矩形（pptx 不原生支持渐变，用两色块叠加近似）
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE['primary']

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(1.5), SLIDE_W, Inches(1.5))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = PALETTE['accent']

    # 标题
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), SLIDE_W - Inches(1.6), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = course_title
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(18)
        r2.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # 底部版本徽章
    badge = slide.shapes.add_textbox(Inches(0.8), SLIDE_H - Inches(1.1), Inches(8), Inches(0.6))
    btf = badge.text_frame
    bp = btf.paragraphs[0]
    br = bp.add_run()
    br.text = '⚡ TeachAny · 互动课件 · PPTX 派生版'
    br.font.size = Pt(14)
    br.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    br.font.bold = True


def add_content_slide(pres, title: str, bullets: list, image_path: Path = None, quizzes: list = None, interactive: list = None):
    """标准内容幻灯片"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank

    # 标题栏
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.9))
    title_bar.line.fill.background()
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = PALETTE['primary']

    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), SLIDE_W - Inches(1), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title or '(无标题)'
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    # 左侧文本区 / 右侧图片区
    has_image = image_path is not None and image_path.exists()
    text_left = Inches(0.5)
    text_top = Inches(1.2)
    text_width = Inches(7.5) if has_image else Inches(12.3)
    text_height = Inches(5.5)

    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    ttf = text_box.text_frame
    ttf.word_wrap = True

    first = True
    for bullet in bullets[:10]:
        p = ttf.paragraphs[0] if first else ttf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = '• ' + bullet
        r.font.size = Pt(16)
        r.font.color.rgb = PALETTE['text_dark']
        r.font.name = 'Microsoft YaHei'

    if has_image:
        try:
            pic_left = Inches(8.2)
            pic_top = Inches(1.2)
            pic_width = Inches(4.7)
            pic_height = Inches(5.5)
            slide.shapes.add_picture(str(image_path), pic_left, pic_top, width=pic_width, height=pic_height)
        except Exception as e:
            print(f"⚠️  插图嵌入失败 ({image_path.name}): {e}")

    # 练习题区（如有）
    if quizzes:
        qb = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), SLIDE_W - Inches(1), Inches(1.1))
        qtf = qb.text_frame
        qtf.word_wrap = True
        p = qtf.paragraphs[0]
        r = p.add_run()
        r.text = '📝 课堂练习（线上版答题）：'
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = PALETTE['accent']
        for q in quizzes[:3]:
            p2 = qtf.add_paragraph()
            r2 = p2.add_run()
            r2.text = q[:120] + ('…' if len(q) > 120 else '')
            r2.font.size = Pt(11)
            r2.font.color.rgb = PALETTE['text_soft']

    # 互动降级提示（脚注）
    if interactive:
        note = slide.shapes.add_textbox(Inches(0.5), SLIDE_H - Inches(0.4), SLIDE_W - Inches(1), Inches(0.3))
        ntf = note.text_frame
        p = ntf.paragraphs[0]
        r = p.add_run()
        r.text = '🔗 含互动组件：' + ' · '.join(interactive) + '（请在 HTML 课件中体验）'
        r.font.size = Pt(10)
        r.font.italic = True
        r.font.color.rgb = PALETTE['text_soft']


def add_interactive_placeholder_slide(pres, section_title: str, components: list, html_link: str):
    """纯互动组件 section 的降级占位页（如整个 section 只是个 canvas 或知识图谱）"""
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE['placeholder']

    # 大标题
    tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), SLIDE_W - Inches(2), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = section_title or '互动环节'
    r.font.size = Pt(36)
    r.font.bold = True
    r.font.color.rgb = PALETTE['primary_dark']

    # 组件列表
    cb = slide.shapes.add_textbox(Inches(1), Inches(3.5), SLIDE_W - Inches(2), Inches(2))
    ctf = cb.text_frame
    ctf.word_wrap = True
    for comp in components:
        p = ctf.paragraphs[0] if ctf.paragraphs[0].text == '' else ctf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = comp
        r.font.size = Pt(22)
        r.font.color.rgb = PALETTE['text_dark']

    # 底部提示
    hb = slide.shapes.add_textbox(Inches(1), SLIDE_H - Inches(1.5), SLIDE_W - Inches(2), Inches(0.8))
    htf = hb.text_frame
    p = htf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '请在 HTML 课件中体验完整互动'
    r.font.size = Pt(14)
    r.font.color.rgb = PALETTE['text_soft']

    if html_link:
        p2 = htf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = html_link
        r2.font.size = Pt(11)
        r2.font.color.rgb = PALETTE['primary']


def add_end_slide(pres, course_title: str):
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = PALETTE['primary_dark']

    tb = slide.shapes.add_textbox(Inches(1), Inches(3), SLIDE_W - Inches(2), Inches(1.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = '谢谢！'
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = RGBColor(0xff, 0xff, 0xff)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f'《{course_title}》 · TeachAny 互动课件'
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xff, 0xff, 0xff)


# ───────────────────────────── 主流程 ─────────────────────────────

def export_pptx(course_dir: Path, output: Path = None) -> Path:
    course_dir = course_dir.resolve()
    html_path = course_dir / 'index.html'
    if not html_path.exists():
        raise FileNotFoundError(f'课件目录缺少 index.html: {html_path}')

    html = html_path.read_text(encoding='utf-8')
    soup = BeautifulSoup(html, 'html.parser')
    manifest = read_manifest(course_dir)

    course_title = (manifest.get('title') or
                    clean_text(soup.title.string if soup.title else '') or
                    course_dir.name)
    subtitle_meta = []
    if manifest.get('subject'):
        subtitle_meta.append(manifest['subject'])
    if manifest.get('grade'):
        subtitle_meta.append(f"G{manifest['grade']}")
    if manifest.get('teachany_version'):
        subtitle_meta.append(f"TeachAny v{manifest['teachany_version']}")
    subtitle = ' · '.join(subtitle_meta)

    pres = setup_pres()
    add_title_slide(pres, course_title, subtitle)

    sections = soup.find_all('section')
    if not sections:
        # 兼容 TeachAny 课件：用 <div class="slide"> 作为幻灯片容器
        sections = soup.select('div.slide')
    if not sections:
        print('⚠️  HTML 中未找到任何 <section> 或 <div class="slide">，仅导出封面与结尾页。')

    slide_count = 1
    for sec in sections:
        sec_id = sec.get('id', '')
        texts = extract_section_texts(sec)
        image = find_first_image(sec, course_dir)
        interactive = detect_interactive_components(sec)

        # 完全空的 section 跳过
        if not texts['title'] and not texts['bullets'] and not texts['quizzes']:
            continue

        # 主要是互动组件、没有实质文本 → 走占位页
        if interactive and not texts['bullets']:
            add_interactive_placeholder_slide(
                pres,
                texts['title'] or sec_id or '互动环节',
                interactive,
                f'{course_dir.name}/index.html#{sec_id}' if sec_id else ''
            )
        else:
            # 如果文本较多，分页（每页最多 10 条 bullet）
            bullets = texts['bullets']
            if len(bullets) <= 10:
                add_content_slide(pres, texts['title'], bullets, image, texts['quizzes'], interactive)
            else:
                for i in range(0, len(bullets), 10):
                    chunk = bullets[i:i + 10]
                    chunk_title = texts['title'] + ('（续）' if i > 0 else '')
                    add_content_slide(
                        pres,
                        chunk_title,
                        chunk,
                        image if i == 0 else None,
                        texts['quizzes'] if i == 0 else None,
                        interactive if i == 0 else None
                    )
        slide_count += 1

    add_end_slide(pres, course_title)

    out_path = output or (course_dir / f"{course_dir.name}.pptx")
    pres.save(str(out_path))
    return out_path


def main():
    parser = argparse.ArgumentParser(
        description='TeachAny · HTML → PPTX 导出工具（v5.34）'
    )
    parser.add_argument('course_dir', help='课件目录，形如 examples/math-m-vertex-form')
    parser.add_argument('-o', '--output', help='输出 pptx 路径，默认 <课件目录>/<目录名>.pptx')
    args = parser.parse_args()

    course_dir = Path(args.course_dir)
    if not course_dir.exists() or not course_dir.is_dir():
        print(f'❌ 课件目录不存在: {course_dir}')
        sys.exit(1)

    output = Path(args.output) if args.output else None
    try:
        out = export_pptx(course_dir, output)
        size = out.stat().st_size
        print(f'✅ PPTX 导出完成: {out}')
        print(f'   大小: {size / 1024:.1f} KB')
        print(f'   幻灯片数: {len(Presentation(str(out)).slides)}')
    except Exception as e:
        print(f'❌ PPTX 导出失败: {e}')
        sys.exit(3)


if __name__ == '__main__':
    main()
